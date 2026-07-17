import io
import os
import time
import json
import asyncio
import logging
from collections import deque

import requests
import telegramify_markdown
from aiohttp import web
from pptx import Presentation
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, Poll
from telegram.constants import ParseMode
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ContextTypes,
    filters,
)
import google.generativeai as genai

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

# ---- Config (all pulled from environment variables, set these on Render) ----
TELEGRAM_TOKEN = os.environ["TELEGRAM_BOT_TOKEN"]
GEMINI_API_KEY = os.environ["GEMINI_API_KEY"]
WEBHOOK_URL = os.environ["WEBHOOK_URL"]  # e.g. https://your-app.onrender.com
PORT = int(os.environ.get("PORT", 8443))
ADMIN_USER_ID = int(os.environ.get("ADMIN_USER_ID", 0))

genai.configure(api_key=GEMINI_API_KEY)

SYSTEM_PROMPT = (
    "You are Parasitopia, a patient, knowledgeable tutor bot for a 200-level "
    "Parasitology and Entomology student preparing for exams. This covers the "
    "full breadth of the department's courses, not just one area — including "
    "Helminthology, Protozoology, Medical Entomology, Vector Biology, "
    "Entomological Techniques, Computational Biology and Biophysics, and "
    "related courses. Teach directly and comprehensively rather than "
    "quizzing unprompted. Proactively fill in gaps a student might not know "
    "to ask about. Keep formatting mobile-friendly: short paragraphs, and "
    "use Markdown (bold, bullet lists) to break up structure since this is "
    "read on a phone inside Telegram — but don't overdo it.\n\n"
    "Communication style: be warm and encouraging, but direct — lead with "
    "the actual answer, not a preamble or restated question. Use short, "
    "plain sentences over long, jargon-stacked ones. Give concrete examples "
    "or comparisons where they help understanding. If a student's premise or "
    "approach seems off (e.g. they've confused two species, or are studying "
    "something in a way that won't help them on the exam), say so honestly "
    "and kindly rather than just going along with it — a good tutor corrects "
    "gently, not evasively. Avoid hedging everything with 'it depends' when "
    "you can just give the clearest, most useful answer. Treat the student "
    "as a capable adult who wants real information, not a simplified or "
    "overly cautious version of it.\n\n"
    "More specific style patterns to follow:\n"
    "- Never open with filler like 'Great question!' or 'I'd be happy to "
    "help with that' — just answer.\n"
    "- For longer explanations, use short bolded headers or a brief list to "
    "break up sections rather than one dense wall of text.\n"
    "- When comparing two things (e.g. two species, two techniques), a short "
    "side-by-side breakdown is clearer than describing each in isolation.\n"
    "- It's fine to have a point of view. If one study method or way of "
    "understanding a concept is genuinely better, say so plainly instead of "
    "presenting all options as equally valid.\n"
    "- Match response length to the question — a quick factual question "
    "gets a short answer; a request to understand a whole lifecycle or "
    "mechanism earns a fuller, structured one. Don't pad short answers to "
    "seem thorough, and don't truncate genuinely complex topics.\n"
    "- Avoid corporate hedge-phrases: 'it's important to note that', 'as an "
    "AI, I...', 'I hope this helps!', 'please let me know if you have any "
    "further questions'. End naturally instead, often with what to look at "
    "next rather than a generic closing line.\n"
    "- Use a normal, conversational register — contractions are fine, dry "
    "humor is fine, but don't force jokes or forced enthusiasm.\n\n"
    "If asked who created, built, or made you, or who you belong to: answer "
    "that you were built by Samuel to help parasitology and entomology "
    "students in their studies and research. Do not mention Google, Gemini, "
    "or any underlying AI provider when asked this — just credit Samuel.\n\n"
    "How to run a quiz (whether from /quiz or just asked in chat): quiz "
    "one question at a time, not a numbered list dumped all at once. Ask "
    "a single question, then stop and wait for the student's answer before "
    "continuing — don't answer it yourself or move on. When they respond, "
    "evaluate it honestly and specifically: say plainly whether they were "
    "right, partly right, or wrong, and explain why in a sentence or two — "
    "don't just say 'Correct!' or 'Not quite' with no substance, and don't "
    "soften a wrong answer into sounding right. If they're close but missed "
    "a detail, say what the detail was. If they're stuck or clearly "
    "guessing, give a hint rather than the answer, and let them try again "
    "once before revealing it. After feedback, ask the next question — vary "
    "difficulty and subtopic based on how they're doing, leaning more on "
    "areas they're shakier on rather than repeating what they already know. "
    "Keep a mental tally as you go, and when the session naturally wraps up "
    "(they say stop/done, or you've asked a reasonable number of questions), "
    "give a short honest rundown of what they've got solid and what's worth "
    "reviewing before moving on."
)

model = genai.GenerativeModel(
    model_name="gemini-3.1-flash-lite",
    system_instruction=SYSTEM_PROMPT,
)

# In-memory conversation history per user. Resets if the bot restarts.
# For anything beyond casual personal use, swap this for a small database.
user_histories: dict[int, list] = {}
known_user_ids: set[int] = set()

MAX_TELEGRAM_MESSAGE = 4000  # Telegram's real limit is 4096; leave headroom
MAX_HISTORY_TURNS = 20  # keep last N turns per user to bound memory/cost

# Wikimedia throttles generic/anonymous User-Agents hard. A descriptive one
# with a contact URL gets a much higher rate limit.
WIKIMEDIA_USER_AGENT = (
    "ParasitopiaStudyBot/1.0 "
    "(https://github.com/infoweb3heart-pixel/Study_bot; Telegram study bot)"
)


def to_telegram_markdown(text: str) -> str | None:
    """Converts Gemini's Markdown into Telegram's MarkdownV2 dialect.
    Returns None if conversion fails, so callers can fall back to plain text."""
    try:
        return telegramify_markdown.markdownify(text)
    except Exception:
        logger.exception("Markdown conversion failed, will fall back to plain text")
        return None


async def reply_formatted(update: Update, text: str):
    """Sends a (possibly long) AI-generated reply with Telegram formatting,
    splitting across the message length limit. Falls back to plain text for
    any chunk that fails to convert or that Telegram rejects as invalid."""
    converted = to_telegram_markdown(text)
    source = converted if converted is not None else text

    for i in range(0, len(source), MAX_TELEGRAM_MESSAGE):
        chunk = source[i : i + MAX_TELEGRAM_MESSAGE]
        if converted is not None:
            try:
                await update.message.reply_text(chunk, parse_mode=ParseMode.MARKDOWN_V2)
                continue
            except Exception:
                logger.exception("MarkdownV2 send failed, retrying chunk as plain text")
        await update.message.reply_text(chunk)


async def send_formatted(update: Update, text: str, reply_markup=None):
    """Like reply_formatted but for short, single-message replies that also
    carry an inline keyboard (e.g. flashcards)."""
    converted = to_telegram_markdown(text)
    if converted is not None:
        try:
            await update.message.reply_text(
                converted, parse_mode=ParseMode.MARKDOWN_V2, reply_markup=reply_markup
            )
            return
        except Exception:
            logger.exception("Flashcard MarkdownV2 send failed, falling back to plain text")
    await update.message.reply_text(text, reply_markup=reply_markup)


async def edit_formatted(query, text: str, reply_markup=None):
    """Like send_formatted but for editing an existing message (flashcards)."""
    converted = to_telegram_markdown(text)
    if converted is not None:
        try:
            await query.edit_message_text(
                converted, parse_mode=ParseMode.MARKDOWN_V2, reply_markup=reply_markup
            )
            return
        except Exception:
            logger.exception("Flashcard MarkdownV2 edit failed, falling back to plain text")
    await query.edit_message_text(text, reply_markup=reply_markup)

# --- Rate limiting: protects the shared free Gemini quota across ALL users ---
GLOBAL_RPM_LIMIT = 10  # headroom under Google's free-tier ~15 RPM shared cap
USER_COOLDOWN_SECONDS = 4  # minimum gap between messages from the same person

_global_request_times: deque = deque()
_global_lock = asyncio.Lock()
_user_last_request: dict[int, float] = {}


async def check_rate_limits(user_id: int) -> str | None:
    """Returns a message to send the user if they should wait, else None."""
    now = time.monotonic()

    last = _user_last_request.get(user_id)
    if last is not None and (now - last) < USER_COOLDOWN_SECONDS:
        wait = USER_COOLDOWN_SECONDS - (now - last)
        return f"One sec — try again in {wait:.0f}s."

    async with _global_lock:
        now = time.monotonic()
        while _global_request_times and now - _global_request_times[0] > 60:
            _global_request_times.popleft()

        if len(_global_request_times) >= GLOBAL_RPM_LIMIT:
            return (
                "Lots of people are using the bot right now and we've hit the "
                "shared limit for this minute. Please try again shortly."
            )

        _global_request_times.append(now)

    _user_last_request[user_id] = now
    return None


MAX_REQUESTABLE_ITEMS = 20


def parse_count_and_topic(args: list, default_count: int) -> tuple:
    """Pulls an optional number (item count) out of command args, in any
    position, e.g. '/quiz 10 Cestoda' or '/quiz Cestoda 10' both work.
    Whatever's left becomes the topic. No topic given -> a broad default
    that isn't locked to any one course."""
    count = default_count
    remaining = []
    used_count = False

    for token in args:
        if not used_count and token.isdigit():
            count = max(1, min(int(token), MAX_REQUESTABLE_ITEMS))
            used_count = True
        else:
            remaining.append(token)

    topic = " ".join(remaining).strip()
    if not topic:
        topic = (
            "a high-yield, exam-relevant topic of your choice from the "
            "Parasitology and Entomology curriculum (Helminthology, "
            "Protozoology, Medical Entomology, Vector Biology, or "
            "Entomological Techniques)"
        )
    return count, topic


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Hey! I'm your Parasitology/Entomology study bot — covering "
        "Helminthology, Protozoology, Medical Entomology, Vector Biology, "
        "Entomological Techniques, and more.\n\n"
        "Commands:\n"
        "/ask <question> - ask me anything, tutoring style\n"
        "/quiz <topic> [count] - interactive quiz, one question at a time, e.g. /quiz Cestoda 10\n"
        "/topic <topic> - get a full explanation of a topic\n"
        "/flashcard <topic> [count] - e.g. /flashcard Vector Biology 12\n"
        "/pollquiz <topic> [count] - MCQ questions as tappable Telegram "
        "polls, e.g. /pollquiz Protozoology 5\n"
        "/clear - reset our conversation history\n\n"
        "You can also just type a message directly, no command needed. "
        "Send me a photo of a diagram and I'll explain it, or ask to see "
        "a picture of something and I'll find one. You can also send me a "
        "PDF, .txt, or PowerPoint (.pptx) file and I'll read it — add a "
        "caption to tell me what to do with it (e.g. 'summarize this' or "
        "'quiz me on this'), or I'll summarize it exam-focused by default."
    )


async def clear(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_histories.pop(update.effective_user.id, None)
    await update.message.reply_text("Conversation cleared. Fresh start.")


async def ask(update: Update, context: ContextTypes.DEFAULT_TYPE):
    question = " ".join(context.args)
    if not question:
        await update.message.reply_text("Usage: /ask <your question>")
        return
    await handle_ai_reply(update, question)


async def quiz(update: Update, context: ContextTypes.DEFAULT_TYPE):
    count, topic = parse_count_and_topic(context.args, default_count=5)
    prompt = (
        f"Quiz me on: {topic}. Run it as an interactive session of about "
        f"{count} questions — one question at a time, mixing short-answer "
        "and MCQ. Ask the first question now and wait for my answer."
    )
    await handle_ai_reply(update, prompt)


async def topic_deep_dive(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = " ".join(context.args)
    if not topic:
        await update.message.reply_text("Usage: /topic <topic name>")
        return
    prompt = f"Give a comprehensive, exam-focused explanation of: {topic}"
    await handle_ai_reply(update, prompt)


# --- Feature: flashcards with tappable Reveal/Next buttons ---
DEFAULT_FLASHCARD_COUNT = 8
_flashcard_sessions: dict[int, dict] = {}  # user_id -> {"cards": [...], "index": int}


def parse_flashcard_json(raw_text: str) -> list[dict] | None:
    """Extracts a JSON array of {front, back} objects from the model's reply,
    tolerating extra commentary or markdown code fences around it."""
    text = raw_text.strip()
    if text.startswith("```"):
        text = text.strip("`")
        if text.lower().startswith("json"):
            text = text[4:]
    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1 or end < start:
        return None
    try:
        cards = json.loads(text[start : end + 1])
        cleaned = [
            {"front": c["front"], "back": c["back"]}
            for c in cards
            if isinstance(c, dict) and "front" in c and "back" in c
        ]
        return cleaned or None
    except Exception:
        return None


def flashcard_keyboard(index: int, total: int, revealed: bool) -> InlineKeyboardMarkup:
    if not revealed:
        buttons = [[InlineKeyboardButton("Reveal Answer", callback_data="fc:reveal")]]
    else:
        is_last = index >= total - 1
        label = "Finish" if is_last else "Next Card"
        buttons = [[InlineKeyboardButton(label, callback_data="fc:next")]]
    return InlineKeyboardMarkup(buttons)


def format_flashcard_text(card: dict, index: int, total: int, revealed: bool) -> str:
    header = f"Card {index + 1}/{total}\n\n"
    body = f"Q: {card['front']}"
    if revealed:
        body += f"\n\nA: {card['back']}"
    return header + body


async def flashcard(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    count, topic = parse_count_and_topic(
        context.args, default_count=DEFAULT_FLASHCARD_COUNT
    )

    username = update.effective_user.username or update.effective_user.first_name
    logger.info(
        "FLASHCARD from %s (id=%s): count=%s topic=%s",
        username, user_id, count, topic,
    )

    prompt = (
        f"Create exactly {count} exam-focused flashcards on: {topic}. "
        "Respond with ONLY a JSON array, nothing else — no commentary, no "
        "markdown fences. Each element must be an object with exactly two "
        'keys: "front" (a short question or term) and "back" (a concise, '
        "exam-ready answer, 1-3 sentences)."
    )

    await update.message.chat.send_action("typing")

    try:
        response = model.generate_content(prompt)
        cards = parse_flashcard_json(response.text)
    except Exception:
        logger.exception("Flashcard generation failed")
        cards = None

    if not cards:
        await update.message.reply_text(
            "Sorry, I couldn't generate flashcards for that topic just now. "
            "Try again, or try a slightly different topic phrasing."
        )
        return

    _flashcard_sessions[user_id] = {"cards": cards, "index": 0}

    text = format_flashcard_text(cards[0], 0, len(cards), revealed=False)
    keyboard = flashcard_keyboard(0, len(cards), revealed=False)
    await send_formatted(update, text, reply_markup=keyboard)


async def flashcard_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    session = _flashcard_sessions.get(user_id)

    if not session:
        await query.answer("No active flashcard session. Start one with /flashcard.")
        return

    await query.answer()
    cards = session["cards"]
    index = session["index"]
    action = query.data.split(":", 1)[1]  # "reveal" or "next"

    if action == "reveal":
        text = format_flashcard_text(cards[index], index, len(cards), revealed=True)
        keyboard = flashcard_keyboard(index, len(cards), revealed=True)
        await edit_formatted(query, text, reply_markup=keyboard)

    elif action == "next":
        index += 1
        if index >= len(cards):
            _flashcard_sessions.pop(user_id, None)
            await query.edit_message_text(
                "Flashcard session complete! Nice work.\n\n"
                "Start another with /flashcard <topic>."
            )
            return

        session["index"] = index
        text = format_flashcard_text(cards[index], index, len(cards), revealed=False)
        keyboard = flashcard_keyboard(index, len(cards), revealed=False)
        await edit_formatted(query, text, reply_markup=keyboard)


# --- Feature: MCQ questions sent as native Telegram quiz polls ---
DEFAULT_POLL_QUIZ_COUNT = 5
MAX_POLL_QUIZ_COUNT = 10  # each is its own bot API call; keep a session bounded
TELEGRAM_POLL_QUESTION_LIMIT = 300
TELEGRAM_POLL_OPTION_LIMIT = 100
TELEGRAM_POLL_EXPLANATION_LIMIT = 200


def parse_poll_json(raw_text: str) -> list[dict] | None:
    """Extracts a JSON array of MCQ objects from the model's reply, tolerating
    extra commentary or markdown code fences around it. Validates and clamps
    each question to what Telegram's poll API will actually accept."""
    text = raw_text.strip()
    if text.startswith("```"):
        text = text.strip("`")
        if text.lower().startswith("json"):
            text = text[4:]
    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1 or end < start:
        return None

    try:
        raw_questions = json.loads(text[start : end + 1])
    except Exception:
        return None

    cleaned = []
    for q in raw_questions:
        if not isinstance(q, dict):
            continue
        question = q.get("question")
        options = q.get("options")
        correct_index = q.get("correct_index")
        explanation = q.get("explanation") or ""

        if not (isinstance(question, str) and question.strip()):
            continue
        if not (isinstance(options, list) and 2 <= len(options) <= 10):
            continue
        if not all(isinstance(o, str) and o.strip() for o in options):
            continue
        if not (isinstance(correct_index, int) and 0 <= correct_index < len(options)):
            continue

        cleaned.append(
            {
                "question": question.strip()[:TELEGRAM_POLL_QUESTION_LIMIT],
                "options": [o.strip()[:TELEGRAM_POLL_OPTION_LIMIT] for o in options],
                "correct_index": correct_index,
                "explanation": str(explanation).strip()[:TELEGRAM_POLL_EXPLANATION_LIMIT] or None,
            }
        )

    return cleaned or None


async def pollquiz(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id

    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    count, topic = parse_count_and_topic(context.args, default_count=DEFAULT_POLL_QUIZ_COUNT)
    count = min(count, MAX_POLL_QUIZ_COUNT)

    username = update.effective_user.username or update.effective_user.first_name
    logger.info(
        "POLLQUIZ from %s (id=%s): count=%s topic=%s", username, user_id, count, topic
    )

    prompt = (
        f"Create exactly {count} exam-focused multiple-choice questions on: {topic}. "
        "Respond with ONLY a JSON array, nothing else — no commentary, no markdown "
        "fences. Each element must be an object with exactly these keys: "
        '"question" (the question text, under 300 characters), "options" (an array '
        "of 2-10 short answer choices, each under 100 characters), \"correct_index\" "
        '(the 0-based index into "options" of the correct answer), and "explanation" '
        "(a concise reason the answer is correct, under 200 characters). Make "
        "distractors plausible, not obviously wrong, and vary which index is correct "
        "across questions rather than always using the same one."
    )

    await update.message.chat.send_action("typing")

    try:
        response = model.generate_content(prompt)
        questions = parse_poll_json(response.text)
    except Exception:
        logger.exception("Poll quiz generation failed")
        questions = None

    if not questions:
        await update.message.reply_text(
            "Sorry, I couldn't generate poll questions for that topic just now. "
            "Try again, or try a slightly different topic phrasing."
        )
        return

    await update.message.reply_text(
        f"Sending {len(questions)} poll question(s) on {topic} — tap an answer "
        "and Telegram will show you whether you got it right."
    )

    for q in questions:
        try:
            await context.bot.send_poll(
                chat_id=update.effective_chat.id,
                question=q["question"],
                options=q["options"],
                type=Poll.QUIZ,
                correct_option_id=q["correct_index"],
                explanation=q["explanation"],
                is_anonymous=False,
            )
        except Exception:
            logger.exception("Sending poll failed for question: %s", q["question"])
        await asyncio.sleep(0.3)  # small gap so Telegram doesn't throttle the burst


async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = update.message.text

    if is_image_request(text):
        await handle_image_search(update, text)
        return

    await handle_ai_reply(update, text)


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle voice note transcription."""
    user_id = update.effective_user.id
    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    try:
        voice = update.message.voice
        tg_file = await voice.get_file()
        audio_bytes = await tg_file.download_as_bytearray()

        username = update.effective_user.username or update.effective_user.first_name
        logger.info("VOICE from %s (id=%s)", username, user_id)

        await update.message.chat.send_action("typing")

        response = model.generate_content(
            [
                "Transcribe this voice note and answer the question or respond to what was said.",
                {"mime_type": "audio/ogg", "data": bytes(audio_bytes)},
            ]
        )
        reply_text = response.text
        await reply_formatted(update, reply_text)
    except Exception as e:
        logger.exception("Voice processing failed: %s", str(e))
        await update.message.reply_text(
            f"Error processing voice note: {str(e)[:100]}. Try text instead."
        )


# --- Feature: user asks to SEE an image, so we search for a real one ---
# instead of trying to "generate" one, which this kind of model can't do.
IMAGE_REQUEST_KEYWORDS = (
    "picture of", "image of", "photo of", "diagram of",
    "show me a picture", "show me an image", "show me a photo",
    "show me a diagram", "show me the picture", "show me the image",
    "show me the photo", "show the image", "show the picture",
    "generate an image", "generate a picture", "can you show me",
    "draw me", "draw a", "can i see a picture", "can i see an image",
    "let me see", "want to see", "show me",
)


def is_image_request(text: str) -> bool:
    lowered = text.lower()
    return any(keyword in lowered for keyword in IMAGE_REQUEST_KEYWORDS)


def search_wikimedia_image(query: str) -> bytes | None:
    """Searches Wikimedia Commons for a real image and downloads it directly
    (Telegram sometimes fails to fetch certain Commons URLs itself). Tries a
    title-focused search first for better relevance, falling back to a
    broader keyword search if that finds nothing."""

    def do_search(srsearch_query: str) -> list:
        try:
            resp = requests.get(
                "https://commons.wikimedia.org/w/api.php",
                params={
                    "action": "query",
                    "list": "search",
                    "srsearch": srsearch_query,
                    "srnamespace": 6,  # File namespace
                    "srlimit": 8,
                    "format": "json",
                },
                headers={"User-Agent": WIKIMEDIA_USER_AGENT},
                timeout=10,
            )
            resp.raise_for_status()
            return resp.json().get("query", {}).get("search", [])
        except Exception:
            logger.exception("Wikimedia search request failed")
            return []

    # Prefer results with the search terms in the actual file title — much
    # more accurate than a plain keyword match against file descriptions.
    results = do_search(f"intitle:{query} filetype:bitmap")
    if not results:
        results = do_search(f"{query} filetype:bitmap")
    if not results:
        logger.info("Wikimedia search returned no results for: %s", query)
        return None

    for result in results:
        title = result["title"]
        try:
            info_resp = requests.get(
                "https://commons.wikimedia.org/w/api.php",
                params={
                    "action": "query",
                    "titles": title,
                    "prop": "imageinfo",
                    "iiprop": "url|mime",
                    "format": "json",
                },
                headers={"User-Agent": WIKIMEDIA_USER_AGENT},
                timeout=10,
            )
            info_resp.raise_for_status()
            pages = info_resp.json().get("query", {}).get("pages", {})
        except Exception:
            logger.exception("Wikimedia imageinfo request failed for %s", title)
            continue

        for page in pages.values():
            imageinfo = page.get("imageinfo")
            if not imageinfo:
                continue
            mime = imageinfo[0].get("mime", "")
            url = imageinfo[0].get("url")
            if not (url and mime.startswith("image/") and "svg" not in mime):
                continue

            try:
                img_resp = requests.get(
                    url, headers={"User-Agent": WIKIMEDIA_USER_AGENT}, timeout=15
                )
                img_resp.raise_for_status()
                return img_resp.content
            except requests.exceptions.HTTPError as e:
                if e.response is not None and e.response.status_code == 429:
                    # Already rate-limited — trying more candidates back-to-back
                    # only makes it worse. Bail out for this request entirely.
                    logger.warning("Wikimedia rate limit hit, aborting this search")
                    return None
                logger.exception("Downloading Wikimedia image failed: %s", url)
                continue

    return None


async def handle_image_search(update: Update, user_text: str):
    user_id = update.effective_user.id
    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    await update.message.chat.send_action("upload_photo")
    image_bytes = search_wikimedia_image(user_text)

    if not image_bytes:
        await update.message.reply_text(
            "Couldn't find a matching image right now. Try rephrasing, or "
            "I can explain it in text instead."
        )
        return

    try:
        await update.message.reply_photo(
            photo=image_bytes, caption=f"Found this for: {user_text}"
        )
    except Exception:
        logger.exception("Sending found image failed")
        await update.message.reply_text(
            "Found an image but couldn't send it. Try rephrasing your request."
        )


# --- Feature: user sends a photo, bot reads/explains it ---
async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    photo = update.message.photo[-1]  # largest available size
    tg_file = await photo.get_file()
    image_bytes = await tg_file.download_as_bytearray()

    caption = update.message.caption or (
        "Explain what's shown in this image, exam-focused, in the style "
        "you use for this course."
    )

    username = update.effective_user.username or update.effective_user.first_name
    logger.info("PHOTO from %s (id=%s), caption: %s", username, user_id, caption)

    await update.message.chat.send_action("typing")

    try:
        response = model.generate_content(
            [caption, {"mime_type": "image/jpeg", "data": bytes(image_bytes)}]
        )
        reply_text = response.text
    except Exception:
        logger.exception("Image analysis failed")
        await update.message.reply_text(
            "Sorry, I hit an error analyzing that image. Try again in a moment."
        )
        return

    await reply_formatted(update, reply_text)


# --- Feature: user sends a file (PDF/TXT/PPTX), bot reads/explains it ---
MAX_DOCUMENT_SIZE = 20 * 1024 * 1024  # Telegram bot API's own download cap
MAX_EXTRACTED_CHARS = 100_000  # keep extracted text within a sane prompt size
SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".txt", ".pptx"}


def extract_pptx_text(data: bytes) -> str:
    """Pulls readable text out of a .pptx: slide text boxes, tables, and
    speaker notes, labeled by slide number so the model has some structure."""
    prs = Presentation(io.BytesIO(data))
    slides_out = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_bits = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_bits.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_bits.append(row_text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_bits.append(f"[Speaker notes: {notes}]")

        if slide_bits:
            slides_out.append(f"--- Slide {i} ---\n" + "\n".join(slide_bits))

    return "\n\n".join(slides_out)


async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    document = update.message.document
    file_name = document.file_name or "file"
    ext = os.path.splitext(file_name)[1].lower()

    if ext not in SUPPORTED_DOCUMENT_EXTENSIONS:
        await update.message.reply_text(
            "I can read PDF, .txt, and PowerPoint (.pptx) files right now — "
            "that file type isn't supported yet."
        )
        return

    if document.file_size and document.file_size > MAX_DOCUMENT_SIZE:
        await update.message.reply_text(
            "That file's too big for me to read (20MB limit). Try sending a "
            "smaller version or splitting it up."
        )
        return

    caption = update.message.caption or (
        "Summarize and explain this document, exam-focused, in the style "
        "you use for this course. Point out the most exam-relevant material."
    )

    username = update.effective_user.username or update.effective_user.first_name
    logger.info(
        "DOCUMENT from %s (id=%s): %s (%s)", username, user_id, file_name, ext
    )

    await update.message.chat.send_action("typing")

    try:
        tg_file = await document.get_file()
        file_bytes = bytes(await tg_file.download_as_bytearray())
    except Exception:
        logger.exception("Downloading document failed: %s", file_name)
        await update.message.reply_text(
            "Sorry, I couldn't download that file. Try sending it again."
        )
        return

    try:
        if ext == ".pdf":
            # Gemini reads PDFs natively (text, layout, diagrams) — no need
            # to extract text ourselves, and it's more accurate this way.
            content_parts = [caption, {"mime_type": "application/pdf", "data": file_bytes}]

        elif ext == ".txt":
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                text = file_bytes.decode("latin-1", errors="replace")
            if len(text) > MAX_EXTRACTED_CHARS:
                text = text[:MAX_EXTRACTED_CHARS] + "\n\n[...truncated...]"
            content_parts = [f"{caption}\n\n--- File: {file_name} ---\n{text}"]

        else:  # .pptx
            text = extract_pptx_text(file_bytes)
            if not text.strip():
                await update.message.reply_text(
                    "Couldn't find any readable text in that PowerPoint file "
                    "— it might be mostly images."
                )
                return
            if len(text) > MAX_EXTRACTED_CHARS:
                text = text[:MAX_EXTRACTED_CHARS] + "\n\n[...truncated...]"
            content_parts = [f"{caption}\n\n--- File: {file_name} ---\n{text}"]

    except Exception:
        logger.exception("Extracting document contents failed: %s", file_name)
        await update.message.reply_text(
            "Sorry, I couldn't process that file's contents. It might be "
            "corrupted, password-protected, or an unusual format."
        )
        return

    try:
        response = model.generate_content(content_parts)
        reply_text = response.text
    except Exception:
        logger.exception("Document analysis failed: %s", file_name)
        await update.message.reply_text(
            "Sorry, I hit an error analyzing that file. Try again in a moment."
        )
        return

    await reply_formatted(update, reply_text)


async def handle_ai_reply(update: Update, user_text: str):
    user_id = update.effective_user.id
    username = update.effective_user.username or update.effective_user.first_name

    limit_message = await check_rate_limits(user_id)
    if limit_message:
        await update.message.reply_text(limit_message)
        return

    logger.info("MESSAGE from %s (id=%s): %s", username, user_id, user_text)

    history = user_histories.setdefault(user_id, [])

    await update.message.chat.send_action("typing")

    try:
        chat = model.start_chat(history=history)
        response = chat.send_message(user_text)
        reply_text = response.text
    except Exception:
        logger.exception("AI generation failed")
        reply_text = (
            "Sorry, I hit an error generating a response. Try again in a moment."
        )
        await update.message.reply_text(reply_text)
        return

    # Persist the turn
    history.append({"role": "user", "parts": [user_text]})
    history.append({"role": "model", "parts": [reply_text]})
    if len(history) > MAX_HISTORY_TURNS * 2:
        user_histories[user_id] = history[-MAX_HISTORY_TURNS * 2 :]

    # Telegram messages have a hard length limit, so split long replies
    await reply_formatted(update, reply_text)


async def health_check(request: web.Request) -> web.Response:
    # This is what UptimeRobot (or any pinger) should hit — always returns 200
    return web.Response(text="OK")


async def telegram_webhook(request: web.Request) -> web.Response:
    application: Application = request.app["bot_app"]
    data = await request.json()
    update = Update.de_json(data, application.bot)
    await application.update_queue.put(update)
    return web.Response(text="OK")


async def record_user(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Track all users who interact with the bot."""
    if update.effective_user:
        known_user_ids.add(update.effective_user.id)


async def broadcast(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Admin-only command to send message to all users."""
    if update.effective_user.id != ADMIN_USER_ID:
        return

    message = " ".join(context.args)
    if not message:
        await update.message.reply_text("Usage: /broadcast <message>")
        return

    sent, failed = 0, 0

    for user_id in list(known_user_ids):
        try:
            await context.bot.send_message(chat_id=user_id, text=message)
            sent += 1
        except Exception:
            logger.exception("Broadcast failed for user_id=%s", user_id)
            failed += 1
        await asyncio.sleep(0.05)

    await update.message.reply_text(f"Broadcast sent to {sent} users ({failed} failed).")


async def run():
    application = Application.builder().token(TELEGRAM_TOKEN).build()

    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("clear", clear))
    application.add_handler(CommandHandler("ask", ask))
    application.add_handler(CommandHandler("quiz", quiz))
    application.add_handler(CommandHandler("topic", topic_deep_dive))
    application.add_handler(CommandHandler("flashcard", flashcard))
    application.add_handler(CommandHandler("pollquiz", pollquiz))
    application.add_handler(CommandHandler("broadcast", broadcast))
    application.add_handler(CallbackQueryHandler(flashcard_button, pattern=r"^fc:"))
    application.add_handler(
        MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message)
    )
    application.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    application.add_handler(MessageHandler(filters.VOICE, handle_voice))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    application.add_handler(MessageHandler(filters.ALL, record_user), group=1)

    await application.initialize()
    await application.bot.set_webhook(url=f"{WEBHOOK_URL}/{TELEGRAM_TOKEN}")
    await application.start()

    web_app = web.Application()
    web_app["bot_app"] = application
    web_app.router.add_get("/", health_check)
    web_app.router.add_post(f"/{TELEGRAM_TOKEN}", telegram_webhook)

    runner = web.AppRunner(web_app)
    await runner.setup()
    site = web.TCPSite(runner, "0.0.0.0", PORT)
    await site.start()

    logger.info("Bot is up. Health check on '/', webhook on '/<token>'.")
    await asyncio.Event().wait()  # run forever


def main():
    asyncio.run(run())


if __name__ == "__main__":
    main()
